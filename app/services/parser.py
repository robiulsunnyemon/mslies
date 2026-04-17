import fitz  # PyMuPDF
from pptx import Presentation
import logging
import os

logger = logging.getLogger(__name__)

def extract_text_from_pdf(filepath: str) -> str:
    text = ""
    try:
        doc = fitz.open(filepath)
        for page in doc:
            text += page.get_text() + "\n"
    except Exception as e:
        logger.error(f"Error parsing PDF {filepath}: {e}")
        raise e
    return text

def extract_text_from_pptx(filepath: str) -> str:
    text = ""
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        logger.error(f"Error parsing PPTX {filepath}: {e}")
        raise e
    return text

def parse_file(filepath: str, file_type: str) -> str:
    if file_type == "PDF":
        return extract_text_from_pdf(filepath)
    elif file_type == "PPTX":
        return extract_text_from_pptx(filepath)
    else:
        raise ValueError("Unsupported file type")
